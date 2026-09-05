"""Local PDF worker. One JSON request on stdin, one JSON result on stdout.

Local documents stay on this Mac. URL conversion explicitly fetches web content.
Each run owns a fresh output directory; input files are opened read-only.
"""
from __future__ import annotations

import difflib
from contextlib import redirect_stdout
import html
import io
import json
import math
import os
from pathlib import Path
import re
import shutil
import signal
import subprocess
import sys
import tempfile
import time
import uuid

import pymupdf as fitz
from PIL import Image, ImageChops, ImageEnhance, ImageOps
from catalog import BY_ID


class UserError(Exception):
    pass


CHILDREN: list[subprocess.Popen] = []


def stop(*_):
    for child in CHILDREN:
        if child.poll() is None:
            try:
                os.killpg(child.pid, signal.SIGTERM)
            except ProcessLookupError:
                pass
    raise UserError("处理已取消。已完成的文件保留在本次输出目录。")


def progress(message, fraction=0):
    print(json.dumps({"message": message, "progress": fraction}, ensure_ascii=False), file=sys.stderr, flush=True)


def run_process(args, timeout=180):
    child = subprocess.Popen([str(a) for a in args], stdout=subprocess.PIPE, stderr=subprocess.PIPE, start_new_session=True)
    CHILDREN.append(child)
    try:
        out, err = child.communicate(timeout=timeout)
        if child.returncode:
            raise UserError("本地转换引擎未能完成处理。请检查源文件是否可正常打开。")
        return out
    except subprocess.TimeoutExpired:
        os.killpg(child.pid, signal.SIGTERM)
        child.communicate()
        raise UserError("本地转换超时。请缩小文件后重试。")
    finally:
        CHILDREN.remove(child)


def find_office():
    paths = [os.environ.get("AIPDF_OFFICE", ""),
             "/Applications/LibreOffice.app/Contents/MacOS/soffice",
             str(Path.home() / ".cache/codex-runtimes/codex-primary-runtime/dependencies/native/libreoffice-headless/libreoffice/LibreOfficeDev.app/Contents/MacOS/soffice")]
    paths += [shutil.which("libreoffice") or "", shutil.which("soffice") or ""]
    return next((p for p in paths if p and Path(p).is_file()), None)


def vision_path():
    paths = [os.environ.get("AIPDF_VISION", ""), str(Path(__file__).resolve().parent / "VisionHelper"),
             str(Path(__file__).resolve().parent.parent / ".build/release/VisionHelper")]
    return next((p for p in paths if p and Path(p).is_file()), None)


def web_path():
    if vision_path():
        sibling=Path(vision_path()).with_name("WebHelper")
        if sibling.is_file(): return str(sibling)
    return None


def health():
    return {"ok": True, "office": bool(find_office()), "vision": bool(vision_path()),
            "engine": fitz.VersionBind}


def number(options, key, default, minimum=None, maximum=None):
    try:
        n = float(options.get(key, default))
    except (ValueError, TypeError):
        raise UserError(f"{key} 必须是数字。")
    if not math.isfinite(n) or (minimum is not None and n < minimum) or (maximum is not None and n > maximum):
        raise UserError(f"{key} 超出允许范围。")
    return n


def flag(options, key, default=False):
    return str(options.get(key, str(default))).lower() == "true"


def parse_pages(spec, total, duplicates=False):
    if not str(spec).strip():
        return list(range(total))
    result = []
    for item in str(spec).replace("，", ",").split(","):
        m = re.fullmatch(r"\s*(\d+)\s*(?:-\s*(\d+)\s*)?", item)
        if not m:
            raise UserError("页码格式不正确，请输入类似 1-3,5 的范围。")
        a, b = int(m[1]), int(m[2] or m[1])
        if not (1 <= a <= total and 1 <= b <= total):
            raise UserError(f"页码超出范围，这份文件共有 {total} 页。")
        result.extend(range(a - 1, b, 1) if a <= b else range(a - 1, b - 2, -1))
    return result if duplicates else list(dict.fromkeys(result))


def open_pdf(path, password=""):
    try:
        doc = fitz.open(path)
    except Exception:
        raise UserError(f"无法读取 {Path(path).name}。文件可能损坏或格式不受支持。")
    if not doc.is_pdf:
        doc.close()
        raise UserError("此工具需要 PDF 文件。")
    if doc.needs_pass and not doc.authenticate(password):
        doc.close()
        raise UserError(f"{Path(path).name} 需要正确的打开密码。")
    if not len(doc):
        doc.close()
        raise UserError("PDF 没有可处理的页面。")
    return doc


def inspect_file(path, password=""):
    with open_pdf(path, password) as doc:
        fields = []
        for page in doc:
            for widget in page.widgets() or []:
                fields.append({"name": widget.field_name, "value": str(widget.field_value or ""),
                               "type": widget.field_type_string, "page": page.number + 1,
                               "choices": widget.choice_values or [],
                               "states": widget.button_states() if widget.field_type in (2, 5) else {}})
        return {"ok": True, "pages": len(doc), "fields": fields}


def save_pdf(doc, path, **kwargs):
    if not len(doc):
        raise UserError("操作后没有剩余页面。")
    # Never append incrementally: redacted/encrypted outputs must not retain old revisions.
    doc.save(path, garbage=4, deflate=True, encryption=fitz.PDF_ENCRYPT_NONE, **kwargs)
    with fitz.open(path) as check:
        if check.page_count != len(doc):
            raise UserError("输出文件验证失败。")


def rgb(value):
    if not re.fullmatch(r"#[0-9A-Fa-f]{6}", value):
        raise UserError("颜色格式应为 #3159D9。")
    return tuple(int(value[i:i+2], 16) / 255 for i in (1, 3, 5))


def area(page, options):
    x, y = number(options, "x", 48, 0), number(options, "y", 72, 0)
    w, h = number(options, "width", 240, 1), number(options, "height", 64, 1)
    # All UI coordinates refer to the visible, rotated page, from its upper left.
    r = fitz.Rect(x, y, x+w, y+h)
    if not page.rect.contains(r):
        raise UserError("选择区域超出了页面，请在预览中重新框选。")
    return r * page.derotation_matrix


def insert_text(page, rect, text, size=18, color=(0, 0, 0), opacity=1, align=0):
    if not text:
        raise UserError("请输入文字内容。")
    if any(ord(c)>255 for c in text):
        css_color="rgb("+",".join(str(round(v*255)) for v in color)+")"
        alignment={0:"left",1:"center",2:"right"}.get(align,"left")
        css=f"* {{font-family:sans-serif;font-size:{size}pt;color:{css_color};text-align:{alignment};}} body {{margin:0;padding:0;}}"
        remaining,_=page.insert_htmlbox(rect,html.escape(text).replace("\n","<br>"),css=css,opacity=opacity,scale_low=1)
    else:
        remaining = page.insert_textbox(rect, text, fontsize=size, fontname="helv", color=color,
                                        fill_opacity=opacity, align=align)
    if remaining < 0:
        raise UserError("文字放不进所选区域，请扩大区域或减小字号。")


def position_rect(page, options, width, height):
    bounds = page.rect
    p = options.get("position", "bottom-center")
    width, height = min(width, bounds.width-32), min(height, bounds.height-32)
    x = 16 if "left" in p else bounds.width-width-16 if "right" in p else (bounds.width-width)/2
    y = 16 if "top" in p else (bounds.height-height)/2 if p == "center" else bounds.height-height-16
    return fitz.Rect(x, y, x+width, y+height) * page.derotation_matrix


def load_image(path, preserve_alpha=False):
    try:
        im = Image.open(path)
        return ImageOps.exif_transpose(im).convert("RGBA" if preserve_alpha else "RGB")
    except Exception:
        helper = vision_path()
        if not helper:
            raise UserError("无法读取图片。请转换为 PNG 或 JPEG。")
        with tempfile.TemporaryDirectory(prefix="aipdf-image-") as tmp:
            out = Path(tmp) / "image.png"
            run_process([helper, "image", path, out])
            return Image.open(out).convert("RGBA" if preserve_alpha else "RGB")


def images_pdf(files, options, dest, scan=False):
    doc = fitz.open()
    for index, path in enumerate(files):
        progress(f"正在导入图片 {index+1}/{len(files)}", index/len(files))
        im = load_image(path)
        if scan and flag(options, "perspective", True):
            helper = vision_path()
            if not helper:
                raise UserError("缺少 Apple Vision 扫描组件，请重新构建应用。")
            with tempfile.TemporaryDirectory(prefix="aipdf-scan-") as tmp:
                source, out = Path(tmp)/"in.png", Path(tmp)/"out.png"
                im.save(source)
                run_process([helper, "scan", source, out])
                im = Image.open(out).convert("RGB")
        if scan and options.get("scanMode") != "original":
            im = ImageOps.autocontrast(ImageOps.grayscale(im))
            if options.get("scanMode", "document") == "document":
                im = ImageEnhance.Contrast(im).enhance(1.5)
        page_size = options.get("pageSize", "A4")
        width, height = (im.width*.75, im.height*.75) if page_size == "image" else fitz.paper_size(page_size)
        if options.get("orientation") == "landscape":
            width, height = height, width
        margin = number(options, "margin", 24, 0)
        if margin*2 >= min(width, height):
            raise UserError("页边距太大，图片没有可用空间。")
        page = doc.new_page(width=width, height=height)
        stream = io.BytesIO()
        im.save(stream, format="PNG")
        page.insert_image(fitz.Rect(margin, margin, width-margin, height-margin), stream=stream.getvalue())
    save_pdf(doc, dest)
    doc.close()
    return [dest]


def office_convert(source, outdir, pdfa=None):
    executable = find_office()
    if not executable:
        raise UserError("此转换需要 LibreOffice。请将 LibreOffice 安装到“应用程序”文件夹后重试。")
    with tempfile.TemporaryDirectory(prefix="aipdf-office-") as tmp:
        profile = Path(tmp) / "profile"
        target = Path(tmp) / "output"
        target.mkdir()
        fmt = "pdf"
        if pdfa:
            version = {"1b": 1, "2b": 2, "3b": 3}[pdfa]
            fmt = 'pdf:draw_pdf_Export:' + json.dumps({"SelectPdfVersion": {"type": "long", "value": str(version)}})
        run_process([executable, f"-env:UserInstallation={profile.as_uri()}", "--headless", "--nologo", "--nodefault", "--norestore",
                     "--convert-to", fmt, "--outdir", target, source], timeout=300)
        output = target / (Path(source).stem + ".pdf")
        if not output.is_file():
            raise UserError("Office 引擎没有生成 PDF。请检查文件是否加密、损坏或无法打开。")
        with fitz.open(output) as check:
            if not len(check):
                raise UserError("转换后的 PDF 为空。")
            if pdfa and f'part="{pdfa[0]}"' not in check.get_xml_metadata() and f'<pdfaid:part>{pdfa[0]}</pdfaid:part>' not in check.get_xml_metadata():
                raise UserError("引擎未输出要求的 PDF/A 标识，已停止导出。")
        shutil.copyfile(output, outdir)
    return [outdir]


def ocr_doc(doc, options):
    helper = vision_path()
    if not helper:
        raise UserError("缺少 Apple Vision 识别组件，请重新构建应用。")
    selected = parse_pages(options.get("pages", ""), len(doc))
    recognized = 0
    with tempfile.TemporaryDirectory(prefix="aipdf-ocr-") as tmp:
        for j, i in enumerate(selected):
            page = doc[i]
            progress(f"本机识别第 {i+1} 页", j/len(selected))
            if page.get_text().strip() and not flag(options, "force"):
                continue
            # Normalize rotated pages so image and OCR coordinates agree.
            if page.rotation:
                page.remove_rotation()
            image = Path(tmp)/"page.png"
            page.get_pixmap(dpi=216).save(image)
            boxes = json.loads(run_process([helper, "ocr", image], timeout=180))
            for item in boxes:
                text = item["text"]
                x, y, w, h = item["x"]*page.rect.width, item["y"]*page.rect.height, item["w"]*page.rect.width, item["h"]*page.rect.height
                font = fitz.Font("china-s" if any(ord(c)>255 for c in text) else "helv")
                size = min(h*.8, w / max(font.text_length(text, fontsize=1), .01))
                page.insert_text((x, y+h*.82), text, fontname=font.name if font.name == "helv" else ("china-s" if any(ord(c)>255 for c in text) else "helv"), fontsize=max(.5, size), render_mode=3)
                recognized += 1
    return recognized


def compress_doc(doc, options):
    quality = options.get("quality", "balanced")
    if quality == "lossless":
        return
    max_dpi, jpeg_quality = (144, 70) if quality == "balanced" else (96, 45)
    seen = set()
    for page in doc:
        for info in page.get_images(full=True):
            xref, smask = info[0], info[1]
            if xref in seen or smask:
                continue
            seen.add(xref)
            try:
                raw = doc.extract_image(xref)
                im = Image.open(io.BytesIO(raw["image"]))
                if im.width < 256 or im.height < 256:
                    continue
                rects = [r for p in doc for r in p.get_image_rects(xref)]
                if not rects:
                    continue
                max_w = max(max(1, int(r.width/72*max_dpi)) for r in rects)
                max_h = max(max(1, int(r.height/72*max_dpi)) for r in rects)
                im.thumbnail((max_w, max_h), Image.Resampling.LANCZOS)
                encoded = io.BytesIO()
                im.convert("RGB").save(encoded, format="JPEG", quality=jpeg_quality, optimize=True)
                if len(encoded.getvalue()) < len(raw["image"]):
                    page.replace_image(xref, stream=encoded.getvalue())
            except (ValueError, OSError):
                continue


def export_word(doc, options, dest):
    from docx import Document
    from docx.shared import Inches
    result = Document()
    visual = options.get("mode") == "visual"
    for index, page in enumerate(doc):
        if index:
            result.add_page_break()
        if visual:
            result.add_picture(io.BytesIO(page.get_pixmap(dpi=144).tobytes("png")), width=Inches(6.2))
            continue
        tables = page.find_tables().tables
        table_bounds = [fitz.Rect(t.bbox) for t in tables]
        events = []
        for block in page.get_text("dict", sort=True)["blocks"]:
            if any(fitz.Rect(block["bbox"]).intersects(r) for r in table_bounds):
                continue
            events.append((block["bbox"][1], "block", block))
        events += [(t.bbox[1], "table", t.extract()) for t in tables]
        for _, kind, item in sorted(events, key=lambda e:e[0]):
            if kind == "table":
                if not item:
                    continue
                table = result.add_table(rows=len(item), cols=max(len(row) for row in item))
                table.style = "Table Grid"
                for r, row in enumerate(item):
                    for c, cell in enumerate(row):
                        table.cell(r, c).text = str(cell or "")
            elif item["type"] == 1:
                result.add_picture(io.BytesIO(item["image"]), width=Inches(min(6.2, (item["bbox"][2]-item["bbox"][0])/72)))
            else:
                for line in item.get("lines", []):
                    p = result.add_paragraph()
                    for span in line["spans"]:
                        run = p.add_run(span["text"])
                        run.bold = bool(span["flags"] & 16)
                        run.italic = bool(span["flags"] & 2)
    result.save(dest)


def export_ppt(doc, options, dest):
    from pptx import Presentation
    from pptx.util import Pt
    presentation = Presentation()
    presentation.slide_width, presentation.slide_height = Pt(doc[0].rect.width), Pt(doc[0].rect.height)
    for page in doc:
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        sx, sy = doc[0].rect.width/page.rect.width, doc[0].rect.height/page.rect.height
        if options.get("mode", "visual") == "visual":
            slide.shapes.add_picture(io.BytesIO(page.get_pixmap(dpi=144).tobytes("png")), 0, 0, presentation.slide_width, presentation.slide_height)
        else:
            for block in page.get_text("dict")["blocks"]:
                x0,y0,x1,y1 = block["bbox"]
                if block["type"] == 1:
                    slide.shapes.add_picture(io.BytesIO(block["image"]), Pt(x0*sx), Pt(y0*sy), Pt((x1-x0)*sx), Pt((y1-y0)*sy))
                else:
                    tf = slide.shapes.add_textbox(Pt(x0*sx), Pt(y0*sy), Pt((x1-x0)*sx+8), Pt((y1-y0)*sy+8)).text_frame
                    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
                    for n,line in enumerate(block.get("lines", [])):
                        para = tf.paragraphs[0] if n == 0 else tf.add_paragraph()
                        for span in line["spans"]:
                            run = para.add_run()
                            run.text = span["text"]
                            run.font.size = Pt(max(1, span["size"]*sy))
                            run.font.bold = bool(span["flags"] & 16)
    presentation.save(dest)


def export_excel(doc, dest):
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment
    book = Workbook()
    book.remove(book.active)
    for page in doc:
        tables = page.find_tables().tables
        data = [t.extract() for t in tables] or [[[line] for line in page.get_text(sort=True).splitlines()] or [["（本页没有可提取文字）"]]]
        for index, rows in enumerate(data):
            sheet = book.create_sheet(f"第{page.number+1}页-表{index+1}")
            for row in rows:
                # Explicit strings prevent spreadsheet formula injection from source documents.
                sheet.append([str(c or "") for c in row])
                for cell in sheet[sheet.max_row]:
                    cell.data_type = "s"
                    cell.alignment = Alignment(vertical="top", wrap_text=True)
            for cell in sheet[1]:
                cell.font = Font(bold=True, color="FFFFFF")
                cell.fill = PatternFill("solid", fgColor="314A72")
            sheet.freeze_panes = "A2"
            for column in sheet.columns:
                sheet.column_dimensions[column[0].column_letter].width = min(70, max(16, max(len(str(c.value or "")) for c in column)+2))
    book.save(dest)


def export_markdown(doc, options, dest):
    sections = []
    assets = dest.parent / (dest.stem + "_assets")
    for i in parse_pages(options.get("pages", ""), len(doc)):
        page = doc[i]
        sections.append(f"<!-- 第 {i+1} 页 -->\n")
        tables = page.find_tables().tables
        boxes = [fitz.Rect(t.bbox) for t in tables]
        events = [(t.bbox[1], "table", t.extract()) for t in tables]
        for block in page.get_text("dict", sort=True)["blocks"]:
            if not any(fitz.Rect(block["bbox"]).intersects(r) for r in boxes):
                events.append((block["bbox"][1], "block", block))
        for j, (_, kind, obj) in enumerate(sorted(events, key=lambda e:e[0])):
            if kind == "table" and obj:
                esc = lambda x: str(x or "").replace("|", "\\|").replace("\n", "<br>")
                sections.append("| " + " | ".join(esc(c) for c in obj[0]) + " |")
                sections.append("| " + " | ".join("---" for _ in obj[0]) + " |")
                sections.extend("| " + " | ".join(esc(c) for c in row) + " |" for row in obj[1:])
            elif kind == "block" and obj["type"] == 1:
                assets.mkdir(exist_ok=True)
                image = assets / f"page-{i+1}-{j}.{obj['ext']}"
                image.write_bytes(obj["image"])
                sections.append(f"![第 {i+1} 页图片]({assets.name}/{image.name})")
            elif kind == "block":
                for line in obj.get("lines", []):
                    text = "".join(s["text"] for s in line["spans"])
                    size = max((s["size"] for s in line["spans"]), default=11)
                    prefix = "# " if size >= 22 else "## " if size >= 17 else "### " if size >= 14 else ""
                    sections.append(prefix + text)
            sections.append("")
        for link in page.get_links():
            if link.get("uri", "").startswith(("https://", "http://", "mailto:")):
                title = page.get_textbox(link["from"]).strip() or "链接"
                sections.append(f"[{title}]({link['uri']})")
    dest.write_text("\n".join(sections), encoding="utf-8")


def compare_pdfs(files, options, out):
    if len(files) != 2:
        raise UserError("比较需要正好两份 PDF，请按旧版、新版的顺序添加。")
    left, right = (open_pdf(p, options.get("password", "")) for p in files)
    result = fitz.open()
    for i in range(max(len(left),len(right))):
        progress(f"正在比较第 {i+1} 页", i/max(len(left),len(right)))
        images = []
        for doc in (left, right):
            if i < len(doc):
                pix = doc[i].get_pixmap(dpi=96, colorspace=fitz.csRGB)
                im = Image.frombytes("RGB", (pix.width,pix.height), pix.samples)
            else:
                im = Image.new("RGB", (600,800), "white")
            im.thumbnail((900,1200))
            canvas = Image.new("RGB", (900,1200), "white")
            canvas.paste(im, (0,0))
            images.append(canvas)
        difference = ImageChops.difference(*images).convert("L").point(lambda v: 150 if v > 25 else 0)
        highlighted = Image.composite(Image.new("RGB",images[1].size,(255,90,95)),images[1],difference)
        page = result.new_page(width=960,height=690)
        page.insert_text((24,24),f"Before - page {i+1}",fontsize=12)
        page.insert_text((500,24),f"After - page {i+1} (changes in red)",fontsize=12)
        for x, im in ((20,images[0]),(500,highlighted)):
            stream=io.BytesIO(); im.save(stream,format="PNG")
            page.insert_image(fitz.Rect(x,40,x+440,670),stream=stream.getvalue())
    pdf = out/"版本对比.pdf"
    save_pdf(result,pdf)
    diff = difflib.HtmlDiff(wrapcolumn=80).make_file(
        "\n".join(p.get_text(sort=True) for p in left).splitlines(),
        "\n".join(p.get_text(sort=True) for p in right).splitlines(),
        fromdesc=html.escape(Path(files[0]).name), todesc=html.escape(Path(files[1]).name), charset="utf-8")
    report=out/"文字差异.html"; report.write_text(diff)
    left.close(); right.close(); result.close()
    return [pdf,report]


def form_update(doc, options):
    if options.get("formMode", "fill") == "fill":
        values = options.get("formValues", {})
        if not values:
            raise UserError("未找到可填写字段。请先读取表单，或选择创建字段。")
        found=set()
        for page in doc:
            for widget in page.widgets() or []:
                if widget.field_name in values:
                    value = values[widget.field_name]
                    if widget.field_type in (fitz.PDF_WIDGET_TYPE_CHECKBOX, fitz.PDF_WIDGET_TYPE_RADIOBUTTON):
                        on = widget.on_state()
                        value = on if str(value).lower() in ("true", "yes", "1", str(on).lower()) else "Off"
                    widget.field_value=value
                    widget.update()
                    found.add(widget.field_name)
        if set(values)-found:
            raise UserError("部分表单字段不存在，请重新读取文件。")
    else:
        pindex = parse_pages(options.get("pages","1"),len(doc))[0]
        page=doc[pindex]
        if page.rotation: page.remove_rotation()
        name=options.get("fieldName", "").strip()
        if not name:
            raise UserError("请输入字段名称。")
        if any(w.field_name==name for p in doc for w in (p.widgets() or [])):
            raise UserError("字段名称已经存在，请使用新名称。")
        kind=options.get("fieldType","text")
        if kind == "radio":
            # pypdf constructs the shared parent and kids for a true radio group.
            return "radio"
        widget=fitz.Widget()
        widget.field_name=name
        widget.field_type={"text":fitz.PDF_WIDGET_TYPE_TEXT,"checkbox":fitz.PDF_WIDGET_TYPE_CHECKBOX,"list":fitz.PDF_WIDGET_TYPE_LISTBOX}[kind]
        widget.rect=area(page,options)
        widget.field_value=options.get("fieldValue","")
        widget.text_fontsize=12
        widget.border_color=(.35,.4,.5)
        widget.border_width=1
        if kind=="checkbox":
            widget.field_value="Yes" if flag(options,"fieldValue") else "Off"
        if kind=="list":
            widget.choice_values=options.get("choices","").splitlines()
            if not widget.choice_values:
                raise UserError("请填写至少一个列表选项。")
        page.add_widget(widget)


def radio_group(dest, options):
    from pypdf import PdfReader, PdfWriter
    from pypdf.generic import DictionaryObject, NameObject, NumberObject, FloatObject, ArrayObject, TextStringObject, DecodedStreamObject
    reader=PdfReader(dest)
    writer=PdfWriter(); writer.clone_document_from_reader(reader)
    page_index=parse_pages(options.get("pages","1"),len(reader.pages))[0]
    page=writer.pages[page_index]
    if page.get("/Rotate",0):
        raise UserError("创建单选字段前请先把该页旋转归零。")
    choices=[c.strip() for c in options.get("choices", "").splitlines() if c.strip()]
    if len(choices)<2:
        raise UserError("单选组至少需要两个选项。")
    x=float(options.get("x",48)); y=float(options.get("y",72)); h=float(options.get("height",64))
    top=float(page.cropbox.top)-y
    left=float(page.cropbox.left)+x
    kids=ArrayObject()
    group=DictionaryObject({NameObject("/FT"):NameObject("/Btn"),NameObject("/T"):TextStringObject(options["fieldName"]),NameObject("/Ff"):NumberObject(1<<15),NameObject("/Kids"):kids,NameObject("/V"):NameObject("/Off")})
    group_ref=writer._add_object(group)
    for i,choice in enumerate(choices):
        side=min(16,h/len(choices))
        y1=top-i*h/len(choices)
        states=DictionaryObject()
        for state in ("Off",f"Option{i+1}"):
            stream=DecodedStreamObject()
            stream.set_data(("0.4 G 0.5 0.5 15 15 re S\n"+("0 g 4 4 8 8 re f" if state!="Off" else "")).encode())
            stream.update({NameObject("/Type"):NameObject("/XObject"),NameObject("/Subtype"):NameObject("/Form"),NameObject("/BBox"):ArrayObject([NumberObject(0),NumberObject(0),NumberObject(16),NumberObject(16)])})
            states[NameObject("/"+state)]=writer._add_object(stream)
        selected=choice==options.get("fieldValue","")
        widget=DictionaryObject({NameObject("/Type"):NameObject("/Annot"),NameObject("/Subtype"):NameObject("/Widget"),NameObject("/Parent"):group_ref,NameObject("/Rect"):ArrayObject([FloatObject(left),FloatObject(y1-side),FloatObject(left+side),FloatObject(y1)]),NameObject("/F"):NumberObject(4),NameObject("/TU"):TextStringObject(choice),NameObject("/AS"):NameObject(f"/Option{i+1}" if selected else "/Off"),NameObject("/AP"):DictionaryObject({NameObject("/N"):states})})
        ref=writer._add_object(widget); kids.append(ref)
        if "/Annots" not in page: page[NameObject("/Annots")]=ArrayObject()
        page["/Annots"].append(ref)
        if selected: group[NameObject("/V")]=NameObject(f"/Option{i+1}")
    if "/AcroForm" not in writer.root_object:
        writer.root_object[NameObject("/AcroForm")]=writer._add_object(DictionaryObject({NameObject("/Fields"):ArrayObject()}))
    writer.root_object["/AcroForm"]["/Fields"].append(group_ref)
    temporary=dest.with_suffix(".radio.pdf")
    writer.write(temporary)
    os.replace(temporary,dest)


def single_pdf(operation, source, options, out, index):
    stem=f"{index:02d}-"+("网页" if source.startswith(("https://","http://")) else Path(source).stem)
    dest=out/f"{stem}-{operation}.pdf"
    if operation in ("word_to_pdf","ppt_to_pdf","excel_to_pdf"):
        return office_convert(source,dest)
    if operation=="html_to_pdf":
        if source.startswith(("https://","http://")):
            helper=web_path()
            if not helper: raise UserError("缺少本地网页渲染组件，请重新构建应用。")
            with tempfile.TemporaryDirectory(prefix="aipdf-web-") as temp:
                raw=Path(temp)/"web.pdf"
                run_process([helper,source,raw],timeout=90)
                with fitz.open(raw) as website:
                    output=fitz.open()
                    for source_page in website:
                        width=source_page.rect.width
                        page_height=width*842/595
                        y=0
                        while y<source_page.rect.height:
                            clip=fitz.Rect(0,y,width,min(y+page_height,source_page.rect.height))
                            page=output.new_page(width=595,height=842)
                            page.show_pdf_page(fitz.Rect(0,0,595,clip.height*595/width),website,source_page.number,clip=clip)
                            y+=page_height
                    save_pdf(output,dest);output.close()
            return [dest]
        markup=Path(source).read_text(encoding="utf-8")
        # Story has no browser/JavaScript/network runtime; archive resolves local assets.
        story=fitz.Story(html=markup,archive=fitz.Archive(str(Path(source).parent)))
        writer=fitz.DocumentWriter(str(dest)); more=1; pages=0
        while more:
            more,_=story.place(fitz.Rect(36,36,559,806)); pages+=1
            if pages>10000: raise UserError("HTML 文档页数超出限制。")
            device=writer.begin_page(fitz.Rect(0,0,595,842)); story.draw(device); writer.end_page()
        writer.close()
        return [dest]
    doc=open_pdf(source,options.get("password",""))
    try:
        selected=parse_pages(options.get("pages",""),len(doc),duplicates=operation=="organize")
        if operation=="split":
            ranges=options.get("ranges", "").replace("；",";")
            groups=[parse_pages(s,len(doc)) for s in ranges.split(";")] if ranges.strip() else [[i] for i in range(len(doc))]
            outputs=[]
            for n,group in enumerate(groups):
                part=fitz.open()
                for p in group: part.insert_pdf(doc,from_page=p,to_page=p)
                path=out/f"{stem}-part-{n+1:03d}.pdf"; save_pdf(part,path); part.close(); outputs.append(path)
            return outputs
        if operation in ("extract","organize","remove"):
            if operation=="remove":
                if not options.get("pages","").strip(): raise UserError("请输入需要删除的页面。")
                selected=[i for i in range(len(doc)) if i not in selected]
            if not selected: raise UserError("操作后没有剩余页面，请至少保留一页。")
            doc.select(selected)
        elif operation=="compress": compress_doc(doc,options)
        elif operation=="repair": pass  # MuPDF recovers broken xrefs on open; full save rebuilds structure.
        elif operation=="ocr":
            recognized=ocr_doc(doc,options)
            progress(f"识别完成，写入 {recognized} 行文字",.95)
        elif operation=="pdfa":
            with tempfile.TemporaryDirectory(prefix="aipdf-pdfa-") as temp:
                unlocked=Path(temp)/"source.pdf"; save_pdf(doc,unlocked)
                return office_convert(unlocked,dest,options.get("pdfaLevel","2b"))
        elif operation=="pdf_to_images":
            outputs=[]
            if options.get("mode")=="embedded":
                seen=set()
                for i in selected:
                    for info in doc[i].get_images(full=True):
                        xref=info[0]
                        if xref in seen: continue
                        seen.add(xref); raw=doc.extract_image(xref)
                        path=out/f"{stem}-image-{xref}.{raw['ext']}"; path.write_bytes(raw["image"]); outputs.append(path)
                if not outputs: raise UserError("所选页面没有可提取的嵌入图片。")
            else:
                for i in selected:
                    path=out/f"{stem}-page-{i+1:03d}.{options.get('format','jpg')}"
                    doc[i].get_pixmap(dpi=int(number(options,"dpi",144,36,600)),colorspace=fitz.csRGB).save(path)
                    outputs.append(path)
            return outputs
        elif operation=="pdf_to_word":
            dest=dest.with_suffix(".docx"); export_word(doc,options,dest); return [dest]
        elif operation=="pdf_to_ppt":
            dest=dest.with_suffix(".pptx"); export_ppt(doc,options,dest); return [dest]
        elif operation=="pdf_to_excel":
            dest=dest.with_suffix(".xlsx"); export_excel(doc,dest); return [dest]
        elif operation=="markdown":
            dest=dest.with_suffix(".md"); export_markdown(doc,options,dest); return [dest]
        elif operation=="protect":
            password=options.get("newPassword","")
            if not password: raise UserError("请设置非空的打开密码。")
            permissions=fitz.PDF_PERM_ACCESSIBILITY
            if flag(options,"allowPrint",True): permissions|=fitz.PDF_PERM_PRINT|fitz.PDF_PERM_PRINT_HQ
            if flag(options,"allowCopy"): permissions|=fitz.PDF_PERM_COPY
            doc.save(dest,garbage=4,deflate=True,encryption=fitz.PDF_ENCRYPT_AES_256,owner_pw=uuid.uuid4().hex,user_pw=password,permissions=permissions)
            with fitz.open(dest) as check:
                if not check.needs_pass or not check.authenticate(password): raise UserError("加密输出验证失败。")
            return [dest]
        elif operation=="unlock": pass
        elif operation=="forms":
            radio=form_update(doc,options)
            save_pdf(doc,dest)
            if radio=="radio": radio_group(dest,options)
            return [dest]
        else:
            hits=0
            for j,i in enumerate(selected):
                page=doc[i]
                if operation in ("edit","sign","numbers","watermark","redact") and page.rotation:
                    page.remove_rotation()
                if operation=="rotate": page.set_rotation((page.rotation+int(options.get("angle",90)))%360)
                elif operation=="crop":
                    rect=area(page,options)
                    offset=page.cropbox_position
                    rect=fitz.Rect(rect.x0+offset.x,rect.y0+offset.y,rect.x1+offset.x,rect.y1+offset.y)
                    page.set_cropbox(rect)
                elif operation=="numbers":
                    text=options.get("format","{n} / {total}").replace("{n}",str(int(number(options,"start",1))+j)).replace("{total}",str(len(selected)))
                    size=number(options,"fontsize",11,1,200)
                    insert_text(page,position_rect(page,options,240,size*2.2),text,size,align=1)
                elif operation=="watermark":
                    opacity=number(options,"opacity",.2,0,1)
                    if options.get("image"):
                        rect=position_rect(page,options,220,100)
                        watermark=load_image(options["image"],preserve_alpha=True)
                        watermark.putalpha(watermark.getchannel("A").point(lambda a:int(a*opacity)))
                        buf=io.BytesIO(); watermark.save(buf,format="PNG")
                        page.insert_image(rect,stream=buf.getvalue())
                    else:
                        size=number(options,"fontsize",42,1,200)
                        rect=position_rect(page,options,page.rect.width-48,size*2.2)
                        insert_text(page,rect,options.get("text",""),size,(.35,.38,.45),opacity,1)
                elif operation in ("edit","sign"):
                    rect=area(page,options)
                    mode=options.get("mode","text")
                    if (operation=="sign" or mode=="ink") and options.get("strokes"):
                        # Signature strokes are normalized to the signature canvas.
                        for stroke in options["strokes"]:
                            points=[fitz.Point(rect.x0+p[0]*rect.width,rect.y0+p[1]*rect.height) for p in stroke]
                            if len(points)>1: page.draw_polyline(points,color=(.1,.15,.25),width=1.6)
                    elif options.get("image") and (operation=="sign" or mode=="image"):
                        image=load_image(options["image"],preserve_alpha=True); buf=io.BytesIO(); image.save(buf,format="PNG")
                        page.insert_image(rect,stream=buf.getvalue())
                    elif mode=="image": raise UserError("请选择要添加的图片。")
                    elif mode=="rectangle": page.draw_rect(rect,color=rgb(options.get("color","#3159D9")),width=2)
                    elif mode=="highlight": page.draw_rect(rect,color=None,fill=(1,.84,.1),fill_opacity=.3)
                    else: insert_text(page,rect,options.get("text",""),number(options,"fontsize",18,1,200),rgb(options.get("color","#3159D9")))
                elif operation=="redact":
                    terms=[t.strip() for t in options.get("text","").splitlines() if t.strip()]
                    rects=[rect for term in terms for rect in page.search_for(term)] if terms else [area(page,options)]
                    for rect in rects: page.add_redact_annot(rect,fill=(0,0,0)); hits+=1
                    if rects: page.apply_redactions(images=2,graphics=2,text=0)
                else: raise UserError("未识别的工具。")
            if operation=="redact":
                if not hits: raise UserError("没有匹配到敏感文字，未生成输出。请检查关键词或改用区域遮盖。")
                doc.scrub(attached_files=True,clean_pages=True,embedded_files=True,hidden_text=True,javascript=True,metadata=True,redactions=True,remove_links=True,reset_fields=True,reset_responses=True,thumbnails=True,xml_metadata=True)
                for page in doc:
                    for annotation in list(page.annots() or []): page.delete_annot(annotation)
                    for widget in list(page.widgets() or []): page.delete_widget(widget)
                doc.set_toc([])
        save_pdf(doc,dest)
        if operation=="compress" and dest.stat().st_size>=Path(source).stat().st_size and not doc.needs_pass:
            shutil.copyfile(source,dest)
        return [dest]
    finally:
        doc.close()


def process(operation, files, options, out):
    if operation in ("images_to_pdf","scan"):
        return images_pdf(files,options,out/("扫描文档.pdf" if operation=="scan" else "图片文档.pdf"),scan=operation=="scan")
    if operation=="compare": return compare_pdfs(files,options,out)
    if operation=="merge":
        if len(files)<2: raise UserError("请至少添加两份 PDF。")
        merged=fitz.open()
        for i,source in enumerate(files):
            progress(f"正在合并 {i+1}/{len(files)}",i/len(files))
            with open_pdf(source,options.get("password","")) as doc:
                merged.insert_pdf(doc,join_duplicates=False)
        dest=out/"合并文档.pdf"; save_pdf(merged,dest); merged.close(); return [dest]
    outputs=[]
    for i,path in enumerate(files):
        progress(f"正在处理文件 {i+1}/{len(files)}",i/len(files))
        outputs.extend(single_pdf(operation,path,options,out,i+1))
    return outputs


def dispatch(request):
    operation=request.get("operation", "")
    if operation=="health": return health()
    if operation=="inspect": return inspect_file(request["files"][0],request.get("options",{}).get("password",""))
    if operation not in BY_ID: raise UserError("未知的工具。")
    files=[str(Path(p).expanduser().resolve()) for p in request.get("files",[])]
    url=request.get("options",{}).get("url","").strip() if operation=="html_to_pdf" else ""
    if not files and not url: raise UserError("请先选择文件。")
    if url and not re.match(r"^https?://[^\s]+$",url): raise UserError("网页地址必须以 https:// 或 http:// 开头。")
    for f in files:
        if not Path(f).is_file(): raise UserError(f"文件不存在：{Path(f).name}")
        if Path(f).suffix.lower().lstrip(".") not in BY_ID[operation]["extensions"]:
            raise UserError(f"{Path(f).name} 的格式不适用于此工具。")
    if not BY_ID[operation]["multiple"] and len(files)>1: raise UserError("此工具一次处理一份文件。")
    options={f["key"]:f["value"] for f in BY_ID[operation]["fields"]}
    options.update(request.get("options",{}))
    if url: files=[url]
    parent=Path(request.get("outputDir","")).expanduser()
    if not parent.is_absolute() or not parent.is_dir(): raise UserError("请选择有效的输出文件夹。")
    out=parent/("AIPDF-"+time.strftime("%Y%m%d-%H%M%S")+"-"+uuid.uuid4().hex[:6])
    out.mkdir(mode=0o700)
    started=time.monotonic()
    try:
        if operation=="workflow":
            steps=options.get("steps",[])
            if not steps: raise UserError("请至少添加一个工作流程步骤。")
            if len(steps)>20: raise UserError("每个流程最多 20 步。")
            current=files
            allowed={"merge","split","extract","remove","organize","compress","repair","ocr","rotate","numbers","watermark","crop","unlock","protect","redact"}
            for i,step in enumerate(steps):
                op=step.get("operation","")
                if op not in allowed: raise UserError("流程中包含不支持的步骤。")
                stepdir=out/f"{i+1:02d}-{op}"; stepdir.mkdir()
                opts={f["key"]:f["value"] for f in BY_ID[op]["fields"]}
                opts.update(step.get("options",{}))
                opts.setdefault("password",options.get("password",""))
                progress(f"步骤 {i+1}/{len(steps)}：{BY_ID[op]['name']}",i/len(steps))
                current=[str(p) for p in process(op,current,opts,stepdir)]
                if op=="protect": options["password"]=opts.get("newPassword","")
            outputs=[Path(p) for p in current]
        else:
            outputs=process(operation,files,options,out)
        if not outputs or any(not p.is_file() or not p.stat().st_size for p in outputs): raise UserError("未生成有效文件。")
        return {"ok":True,"outputs":[str(p) for p in outputs],"outputDir":str(out),
                "inputBytes":sum(Path(p).stat().st_size for p in files if not p.startswith(("http://","https://"))),"outputBytes":sum(p.stat().st_size for p in outputs),
                "elapsed":round(time.monotonic()-started,2)}
    except Exception as error:
        # Retain any completed results, but do not label the job successful.
        raise UserError(f"{str(error) if isinstance(error,UserError) else '处理失败，请检查文件格式和参数。'}\n本次输出目录：{out}") from error


def main():
    signal.signal(signal.SIGTERM,stop)
    signal.signal(signal.SIGINT,stop)
    try:
        request=json.load(sys.stdin)
        # Third-party parsers may print notices; keep stdout strictly machine-readable.
        with redirect_stdout(sys.stderr):
            result=dispatch(request)
    except (UserError,KeyError,ValueError,TypeError) as error:
        result={"ok":False,"error":str(error) if isinstance(error,UserError) else "请求格式不正确。"}
    except Exception:
        result={"ok":False,"error":"处理引擎发生错误。请检查文件是否正常并重新尝试。"}
    print(json.dumps(result,ensure_ascii=False),flush=True)
    return 0 if result.get("ok") else 1


if __name__=="__main__":
    sys.exit(main())
